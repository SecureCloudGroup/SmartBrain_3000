"""Word / PowerPoint / Excel ingest — with REAL files, through the real parsers.

These formats were explicitly rejected at ingest ("Word/Office ... aren't supported"), which is a
hard wall for a knowledge base: most people's documents are .docx. The files here are built with the
same libraries the app parses them with, so the whole path is exercised end to end — no mocks.

Sections are cited by their real name: a PowerPoint has SLIDES and a spreadsheet has SHEETS, so a
citation says "slide 5", not "p.5". Word has no page map at all — its pagination is decided by the
renderer, so the file genuinely doesn't know where its pages break, and we say so instead of
inventing page numbers.
"""

from __future__ import annotations

import io

import duckdb
import pytest

from smartbrain_3000 import db as dbmod
from smartbrain_3000 import ingest
from smartbrain_3000.kb import KnowledgeBase
from smartbrain_3000.secrets import gen_master_key


def _kb() -> KnowledgeBase:
    conn = duckdb.connect(":memory:")
    dbmod.run_migrations(conn)
    return KnowledgeBase(conn, gen_master_key())


def make_docx(paragraphs: list[str], table: list[list[str]] | None = None, title: str = "") -> bytes:
    import docx

    d = docx.Document()
    if title:
        d.core_properties.title = title
    for p in paragraphs:
        d.add_paragraph(p)
    if table:
        t = d.add_table(rows=len(table), cols=len(table[0]))
        for r, row in enumerate(table):
            for c, val in enumerate(row):
                t.cell(r, c).text = val
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def make_pptx(slides: list[str]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for text in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
        box.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def make_xlsx(sheets: dict[str, list[list[str]]]) -> bytes:
    import openpyxl

    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    for name, rows in sheets.items():
        ws = wb.create_sheet(title=name)
        for row in rows:
            ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# --- Word ------------------------------------------------------------------------------------

def test_docx_text_and_tables_are_extracted() -> None:
    data = make_docx(
        ["The engagement begins on July 10.", "Fees are set out below."],
        table=[["Phase", "Fee"], ["Offering", "$275,000"]],
    )
    title, text, meta = ingest.from_file("Engagement.docx", data)
    assert "engagement begins" in text
    assert "$275,000" in text, "table cells carry real content (fee schedules) — they must be read"
    assert title == "Engagement.docx"  # no embedded title -> fall back to the filename
    assert meta["filename"] == "Engagement.docx"


def test_docx_reports_no_page_map_rather_than_inventing_one() -> None:
    _, _, meta = ingest.from_file("Notes.docx", make_docx(["a", "b"]))
    assert meta["pages"] == [], "Word pagination is a renderer decision; the file doesn't know it"
    assert meta["page_label"] == ""


def test_docx_keeps_a_real_embedded_title() -> None:
    title, _, _ = ingest.from_file("download (3).docx", make_docx(["body"], title="Q3 Earnings Report"))
    assert title == "Q3 Earnings Report"


def test_corrupt_docx_is_a_clean_ingest_error_not_a_500() -> None:
    with pytest.raises(ingest.IngestError, match="Word"):
        ingest.from_file("broken.docx", b"PK\x03\x04 not really a docx")


# --- PowerPoint ------------------------------------------------------------------------------

def test_pptx_slides_become_citable_sections() -> None:
    data = make_pptx(["Opening remarks", "Market overview", "The QUOKKA strategy"])
    _, text, meta = ingest.from_file("Deck.pptx", data)
    assert "QUOKKA" in text
    assert len(meta["pages"]) == 3, "one section per slide"
    assert meta["page_label"] == "slide"
    assert text[meta["pages"][2]:].startswith("The QUOKKA")


def test_a_pptx_hit_is_cited_as_a_slide_not_a_page() -> None:
    kb = _kb()
    title, text, meta = ingest.from_file("Deck.pptx", make_pptx(["intro", "middle", "the QUOKKA strategy"]))
    kb.add(title, text, meta)
    hit = kb.search("quokka")[0]
    assert hit["page"] == 3
    assert hit["page_label"] == "slide"  # the UI renders "slide 3", not "p.3"
    assert hit["source"] == "Deck.pptx"


# --- Excel -----------------------------------------------------------------------------------

def test_xlsx_sheets_become_citable_sections() -> None:
    data = make_xlsx({"Summary": [["Total", "100"]], "Detail": [["Item", "WOMBAT"], ["Qty", "3"]]})
    _, text, meta = ingest.from_file("Book.xlsx", data)
    assert "WOMBAT" in text and "Total" in text
    assert len(meta["pages"]) == 2, "one section per sheet"
    assert meta["page_label"] == "sheet"


def test_an_xlsx_hit_is_cited_as_a_sheet() -> None:
    kb = _kb()
    title, text, meta = ingest.from_file("Book.xlsx", make_xlsx({"A": [["x"]], "B": [["WOMBAT"]]}))
    kb.add(title, text, meta)
    hit = kb.search("wombat")[0]
    assert hit["page"] == 2 and hit["page_label"] == "sheet"


# --- dispatch + the old rejection ------------------------------------------------------------

def test_office_files_are_no_longer_rejected() -> None:
    # The whole point: this used to raise "Word/Office ... aren't supported".
    for name, data in [
        ("a.docx", make_docx(["hello"])),
        ("b.pptx", make_pptx(["hello"])),
        ("c.xlsx", make_xlsx({"S": [["hello"]]})),
    ]:
        _, text, _ = ingest.from_file(name, data)
        assert "hello" in text


def test_office_dispatch_by_content_type_from_a_url(monkeypatch) -> None:
    # Office files are all ZIPs, so magic bytes can't tell them apart — content-type must.
    ctype = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    monkeypatch.setattr(
        ingest.netguard, "safe_fetch_bytes",
        lambda url: {"content": make_docx(["fetched word content"]), "content_type": ctype,
                     "final_url": "https://example.com/report.docx"},
    )
    _, text, meta = ingest.from_url("https://example.com/report.docx")
    assert "fetched word content" in text
    assert meta["source_url"] == "https://example.com/report.docx"


def test_an_unknown_binary_is_still_rejected_with_the_updated_message() -> None:
    with pytest.raises(ingest.IngestError, match="Images and other binaries"):
        ingest._dispatch(b"\x00\x01\x02\xff\xfe binary junk", "", "")
