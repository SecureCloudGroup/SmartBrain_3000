"""Extract text from a URL or uploaded file and store it in the Knowledge base.

URL fetches go through the SSRF guard (netguard.safe_fetch_bytes). PDFs are parsed with pypdf,
Word/PowerPoint/Excel with python-docx / python-pptx / openpyxl, HTML with trafilatura, and
text-family files are decoded directly. Stored documents are embedded best-effort on add so they
are immediately searchable. Heavy parsers are imported lazily (only on the ingest path) to keep app
startup fast.

Formats that HAVE natural sections record where each one starts, so a search hit can be cited to it:
PDF pages, PowerPoint slides, Excel sheets. Word does NOT: its pagination is decided by the renderer
(fonts, margins, printer), so the file itself doesn't know where its pages break — we report no page
map rather than invent page numbers.
"""

from __future__ import annotations

import io
import logging
import os
from collections.abc import Iterable
from urllib.parse import urlparse

import httpx

from . import gateway, kb, netguard

log = logging.getLogger(__name__)

_MAX_TEXT = 1_000_000  # cap on stored text per document (chars)
_MAX_SECTIONS = 1000  # bounded loop over pages / slides / sheets
_MAX_PARTS = 100_000  # bounded loop over paragraphs + table rows in a Word document
_MAX_ROWS = 10_000  # bounded rows per spreadsheet sheet (a sheet can claim a million)
_EMBED_TIMEOUT = 15.0  # interactive embed-on-add / rename: fail fast (best-effort; backfilled later)
# A cold local embed model (e.g. MLX/oMLX loading bge-m3) can take ~50s on its FIRST request.
# The backfill path waits it out so a single reindex succeeds instead of timing out at 15s (then
# working on the 2nd try). Bifrost's per-provider timeout (300s) is the real upstream ceiling.
_REINDEX_EMBED_TIMEOUT = 120.0
_TEXT_EXT = frozenset({".txt", ".md", ".markdown", ".csv", ".json", ".log", ".rst", ".text"})
_HTML_EXT = frozenset({".html", ".htm"})
# Document extensions that mark an extracted title as "filename-shaped" rather than a real title —
# e.g. a PDF exported from Word whose embedded /Title metadata is still the original "Report.DOCX".
# When the extracted title ends in one of these we distrust it and use the real uploaded filename
# / URL tail instead, so the KB never shows a wrong/stale extension (like .DOCX on a .pdf).
_FILENAME_EXTS = frozenset({
    ".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".rtf", ".odt", ".ods", ".odp",
    ".pages", ".key", ".numbers", ".epub", ".txt", ".md", ".markdown", ".html", ".htm", ".csv", ".json",
})


class IngestError(Exception):
    """Content could not be fetched or extracted into usable text."""


def _resolve_title(extracted: str, fallback: str) -> str:
    """Prefer the extractor's title, but use ``fallback`` (the uploaded filename / URL tail) when
    the extracted title is empty OR looks like a filename ending in a document extension — e.g. a
    PDF whose /Title metadata is still the original 'Report.DOCX'. This keeps a wrong/stale
    extension out of the KB name. Splitext removes only the last segment, so a real title like
    'v1.2 spec' (ext '.2 spec', not a doc ext) is kept."""
    t = extracted.strip()
    _, ext = os.path.splitext(t)
    return t if (t and ext.lower() not in _FILENAME_EXTS) else fallback


def _title_from_url(url: str) -> str:
    """Derive a human-ish title from a URL (last path segment, else host)."""
    assert isinstance(url, str), "url must be str"
    assert url, "url required"
    parsed = urlparse(url)
    tail = parsed.path.rsplit("/", 1)[-1] if parsed.path else ""
    return tail or parsed.netloc or url


def _looks_binary(data: bytes) -> bool:
    """True if the bytes look binary (a NUL byte in the first chunk) — unsafe as text."""
    assert isinstance(data, (bytes, bytearray)), "data must be bytes"
    return b"\x00" in data[:1024]


def _bounded_sections(pieces: Iterable[str], limit: int) -> list[str]:
    """Take section texts until the cumulative _MAX_TEXT cap, bounded by ``limit`` sections.

    Capping DURING extraction (not after) is what stops a single text-bomb page/slide/sheet from
    blowing memory before _MAX_TEXT can be applied.
    """
    assert limit >= 1, "limit must be positive"
    out: list[str] = []
    total = 0
    for i, piece in enumerate(pieces):
        if i >= limit:
            break
        if total + len(piece) >= _MAX_TEXT:
            out.append(piece[: _MAX_TEXT - total])
            break
        out.append(piece)
        total += len(piece)
    return out


def _join_sections(sections: list[str]) -> tuple[str, list[int]]:
    """Join per-page/slide/sheet text and record where each section STARTS in the result.

    This is the one and only place the offset arithmetic lives, because it is the arithmetic a
    citation depends on: get it wrong and "page 12" points at page 11. Offsets stay exact across the
    single separator char that join inserts, and across the leading strip.
    """
    starts: list[int] = []
    pos = 0
    for piece in sections:  # bounded by the caller
        starts.append(pos)
        pos += len(piece) + 1  # +1 for the "\n" the join inserts below
    text = "\n".join(sections)
    lead = len(text) - len(text.lstrip())  # the strip shifts every offset left by this much
    return text.strip(), [max(0, p - lead) for p in starts]


def _extract_pdf(data: bytes) -> tuple[str, str, list[int]]:
    """Return (title, text, page_starts) from PDF bytes via pypdf (bounded pages + cumulative cap).

    ``page_starts[i]`` is the character offset in ``text`` where page i+1 begins — the page
    boundaries are known HERE and nowhere else, and used to be flattened away by the join.

    All pypdf access (incl. metadata) is inside the try so any parser error surfaces as IngestError,
    not a 500.
    """
    from pypdf import PdfReader  # lazy: heavy dep, only on the ingest path

    assert isinstance(data, (bytes, bytearray)), "data must be bytes"
    assert data, "data is empty"
    try:
        reader = PdfReader(io.BytesIO(data))
        sections = _bounded_sections((p.extract_text() or "" for p in reader.pages), _MAX_SECTIONS)
        title = (reader.metadata.title if reader.metadata and reader.metadata.title else "") or ""
    except Exception as exc:  # malformed / encrypted PDF
        raise IngestError(f"could not read PDF: {exc}") from None
    text, starts = _join_sections(sections)
    return title.strip(), text, starts


def _extract_docx(data: bytes) -> tuple[str, str, list[int]]:
    """Return (title, text, []) from a Word .docx — paragraphs plus table cells.

    No page map: Word pagination is decided by the renderer (fonts, margins, printer), so the file
    itself does not know where its pages break. We say so honestly rather than invent page numbers.
    """
    import docx  # lazy: only on the ingest path

    assert data, "data is empty"
    try:
        doc = docx.Document(io.BytesIO(data))
        parts = [p.text for p in doc.paragraphs]
        for table in doc.tables:  # tables carry real content (fee schedules, terms) — don't skip them
            for row in table.rows:
                parts.append("\t".join(c.text for c in row.cells))
        title = (doc.core_properties.title or "") if doc.core_properties else ""
    except Exception as exc:  # not a real docx / corrupt zip
        raise IngestError(f"could not read Word document: {exc}") from None
    text, _ = _join_sections(_bounded_sections(parts, _MAX_PARTS))
    return title.strip(), text, []


def _extract_pptx(data: bytes) -> tuple[str, str, list[int]]:
    """Return (title, text, slide_starts) from a PowerPoint .pptx — one section per SLIDE."""
    from pptx import Presentation  # lazy

    assert data, "data is empty"
    try:
        prs = Presentation(io.BytesIO(data))
        slides = [
            "\n".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
            for slide in prs.slides
        ]
        title = (prs.core_properties.title or "") if prs.core_properties else ""
    except Exception as exc:
        raise IngestError(f"could not read PowerPoint file: {exc}") from None
    text, starts = _join_sections(_bounded_sections(slides, _MAX_SECTIONS))
    return title.strip(), text, starts  # a slide is this format's "page" — cited as "slide 5"


def _extract_xlsx(data: bytes) -> tuple[str, str, list[int]]:
    """Return (title, text, sheet_starts) from an Excel .xlsx — one section per SHEET, rows as TSV."""
    import openpyxl  # lazy

    assert data, "data is empty"
    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sheets: list[str] = []
        for ws in wb.worksheets[:_MAX_SECTIONS]:  # bounded
            rows: list[str] = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= _MAX_ROWS:  # bounded: a spreadsheet can claim a million rows
                    break
                rows.append("\t".join("" if v is None else str(v) for v in row))
            sheets.append(f"{ws.title}\n" + "\n".join(rows))
        wb.close()
    except Exception as exc:
        raise IngestError(f"could not read Excel file: {exc}") from None
    text, starts = _join_sections(_bounded_sections(sheets, _MAX_SECTIONS))
    return "", text, starts  # a sheet is this format's "page" — cited as "sheet 2"


def _extract_html(html: str, url: str) -> tuple[str, str]:
    """Return (title, main-article-text) from HTML via trafilatura."""
    import trafilatura  # lazy: heavy dep

    assert isinstance(html, str), "html must be str"
    assert isinstance(url, str), "url must be str"  # may be empty (file upload)
    try:
        text = trafilatura.extract(html, url=url or None, include_comments=False) or ""
        meta = trafilatura.extract_metadata(html)
    except Exception as exc:  # malformed / pathological HTML (lxml errors, recursion)
        raise IngestError(f"could not read HTML: {exc}") from None
    title = meta.title if (meta is not None and getattr(meta, "title", None)) else ""
    return (title or "").strip(), text.strip()


def _clamp_pages(page_starts: list[int], text_len: int) -> list[int]:
    """Drop page starts that fall past the (possibly truncated) text, so no citation points nowhere."""
    return [p for p in page_starts if p < text_len]


# Extension -> (extractor, what one section is called in a citation). Office files are ZIPs, so they
# all share the same magic bytes and can only be told apart by extension / content-type.
_OFFICE = {
    ".docx": (_extract_docx, ""),
    ".pptx": (_extract_pptx, "slide"),
    ".xlsx": (_extract_xlsx, "sheet"),
}
_OFFICE_CTYPE = {
    "wordprocessingml": (_extract_docx, ""),
    "presentationml": (_extract_pptx, "slide"),
    "spreadsheetml": (_extract_xlsx, "sheet"),
}


def _dispatch(data: bytes, content_type: str, hint_url: str) -> tuple[str, str, list[int], str]:
    """Pick an extractor by magic bytes / content-type, falling back to the URL hint.

    Returns (title, text, section_starts, page_label). ``section_starts`` is empty for formats with
    no natural sections, and ``page_label`` names what a section IS ("page", "slide", "sheet").

    Magic bytes and content-type are authoritative; the ``.pdf`` URL hint is only consulted when the
    content-type is generic (so an HTML page served at a ``.pdf`` URL is never handed to the PDF
    parser).
    """
    assert isinstance(data, (bytes, bytearray)), "data must be bytes"
    assert data, "data is empty"
    ctl = content_type.lower()
    if data[:5] == b"%PDF-" or "pdf" in ctl:
        return (*_extract_pdf(data), "page")
    for marker, (extract, label) in _OFFICE_CTYPE.items():
        if marker in ctl:
            return (*extract(data), label)
    if "html" in ctl or "xml" in ctl:
        return (*_extract_html(data.decode("utf-8", "replace"), hint_url), [], "")
    generic = ctl.startswith("application/octet-stream") or not ctl
    if generic and hint_url.lower().endswith(".pdf"):
        return (*_extract_pdf(data), "page")
    ext = os.path.splitext(urlparse(hint_url).path)[1].lower()
    if generic and ext in _OFFICE:
        extract, label = _OFFICE[ext]
        return (*extract(data), label)
    if hint_url.lower().endswith((".html", ".htm")):
        return (*_extract_html(data.decode("utf-8", "replace"), hint_url), [], "")
    if _looks_binary(data):  # don't store an unrecognized binary blob as garbled text
        raise IngestError(
            "unsupported file type — supported: PDF, Word/PowerPoint/Excel (.docx, .pptx, .xlsx), "
            "HTML, and text (.txt, .md, .csv, .json). Images and other binaries aren't supported."
        )
    return "", data.decode("utf-8", "replace").strip(), [], ""  # text/* or json — store as-is


def from_url(url: str) -> tuple[str, str, dict]:
    """Fetch a URL (SSRF-guarded) and extract (title, text, meta). Raises IngestError."""
    assert url, "url required"
    try:
        got = netguard.safe_fetch_bytes(url)
    except netguard.FetchError as exc:
        raise IngestError(str(exc)) from None
    title, text, pages, label = _dispatch(got["content"], got["content_type"], got["final_url"])
    if not text:
        raise IngestError("no readable text found at that URL")
    text = text[:_MAX_TEXT]
    meta = {
        "source_url": got["final_url"],
        "mime": got["content_type"],
        "pages": _clamp_pages(pages, len(text)),
        "page_label": label,
    }
    return _resolve_title(title, _title_from_url(got["final_url"])), text, meta


def from_file(filename: str, data: bytes) -> tuple[str, str, dict]:
    """Extract (title, text, meta) from uploaded file bytes by extension/sniff."""
    assert filename, "filename required"
    assert data, "file is empty"
    ext = os.path.splitext(filename)[1].lower()
    pages: list[int] = []
    label = ""
    if ext == ".pdf" or data[:5] == b"%PDF-":
        title, text, pages = _extract_pdf(data)
        label = "page"
    elif ext in _OFFICE:  # Office files are ZIPs — only the extension tells them apart
        extract, label = _OFFICE[ext]
        title, text, pages = extract(data)
    elif ext in _HTML_EXT:
        title, text = _extract_html(data.decode("utf-8", "replace"), "")
    elif ext in _TEXT_EXT or ext == "":
        if _looks_binary(data):  # an extension-less/text-typed file that is really binary
            raise IngestError("binary file without a recognized text type")
        title, text = "", data.decode("utf-8", "replace").strip()
    else:
        raise IngestError(f"unsupported file type: {ext or 'unknown'}")
    if not text:
        raise IngestError("no readable text found in that file")
    text = text[:_MAX_TEXT]
    # The filename is what a citation actually shows the user ("Lease.pdf, p.12"), and it used to be
    # consumed to derive a title and then dropped.
    meta = {
        "filename": os.path.basename(filename),
        "pages": _clamp_pages(pages, len(text)),
        "page_label": label,
    }
    return _resolve_title(title, os.path.basename(filename)), text, meta


def embed_doc(knowledge, doc_id: str, title: str, content: str, model: str, *, timeout: float = _EMBED_TIMEOUT) -> None:
    """Chunk title+content, embed each chunk, and store the per-chunk vectors so a
    long document is fully searchable (not just its head). ``timeout`` is per embed request —
    the backfill path raises it so a cold local model's first load doesn't fail the embed."""
    assert doc_id and title and model, "doc id, title, model required"
    assert knowledge is not None, "knowledge base required"
    chunks = kb.chunk_text(title, content)  # bounded by _MAX_CHUNKS
    with httpx.Client(base_url=gateway.gateway_url(), timeout=timeout) as client:
        vectors = [gateway.embed(c, model, client=client, timeout=timeout) for c in chunks]
    knowledge.put_embeddings(doc_id, vectors, model)


def reindex_pending(knowledge, model: str, *, limit: int = 1000, timeout: float = _REINDEX_EMBED_TIMEOUT) -> tuple[int, int, int, str]:
    """Backfill embeddings for docs missing one or on a stale model (best-effort, bounded).

    Uses a generous per-embed ``timeout`` so a COLD local embed model (which can take ~50s to
    load on its first request) is waited out rather than cut short — otherwise the first reindex
    fails and only the second (warm) one works. Returns (embedded, skipped, failed, first_error).
    Shared by the manual /api/kb/reindex route and the scheduler's automated backfill."""
    assert knowledge is not None and model, "knowledge + model required"
    pending = knowledge.docs_needing_embedding(model)[:limit]
    embedded = skipped = failed = 0
    error = ""
    for doc_id in pending:  # bounded by limit
        doc = knowledge.get(doc_id)
        if doc is None:
            skipped += 1
            continue
        try:
            embed_doc(knowledge, doc_id, doc["title"], doc["content"], model, timeout=timeout)
            embedded += 1
        except Exception as exc:  # keep going on a single failure
            failed += 1
            if not error:
                error = str(getattr(exc, "message", exc))
    return embedded, skipped, failed, error


def store(knowledge, title: str, content: str, meta: dict | None = None) -> dict:
    """Add a document to the KB (with provenance ``meta``) and embed it best-effort.

    Embedding failure (e.g. Ollama down) never fails the add — the doc is stored
    and a later reindex backfills the vector.
    """
    assert knowledge is not None, "knowledge base required"
    assert title and content, "title + content required"
    doc_id = knowledge.add(title, content, meta)
    try:  # best-effort: make it immediately semantic-searchable
        model = gateway.embed_model(getattr(knowledge, "conn", None))  # routed embedding model
        embed_doc(knowledge, doc_id, title, content, model)
    except Exception as exc:  # embeddings optional; /api/kb/reindex can backfill
        log.info("embed-on-add skipped for %s: %s", doc_id, exc)
    return {"id": doc_id, "title": title, "chars": len(content)}


def ingest_url(knowledge, url: str) -> dict:
    """Fetch + extract + store a URL. Returns {id, title, chars}."""
    title, text, meta = from_url(url)
    return store(knowledge, title, text, meta)
